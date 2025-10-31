from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
import colorsys

class SlideAnalyzer:

    def __init__(self, filepath, fileName):
        self.prs = Presentation(filepath)
        self.slides_data = []
        self.fileName = fileName
        


    def parse_slides(self):
        """Extract all slide data"""
        for idx, slide in enumerate(self.prs.slides):
            slide_info = {
                'slide_number': idx + 1,
                'text_elements': [],
                'images': [],
                'has_title': False,
                'background_color': None
            }
            
            for shape in slide.shapes:
                # Text extraction
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_data = self._extract_text_properties(shape, paragraph)
                            slide_info['text_elements'].append(text_data)
                            
                            # Check if title
                            if shape.is_placeholder:
                                try:
                                    if shape.placeholder_format.type == 1:  # Title
                                        slide_info['has_title'] = True
                                except:
                                    pass
                
                # Image detection
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info['images'].append({
                        'width': shape.width,
                        'height': shape.height
                    })
            
            self.slides_data.append(slide_info)
        
        return self.slides_data
    


    def _extract_text_properties(self, shape, paragraph):
        """Extract font size, color, and text"""
        text = paragraph.text.strip()
        
        # Get font size (handle None)
        font_size = None
        if paragraph.runs:
            run = paragraph.runs[0]
            if run.font.size:
                font_size = run.font.size.pt
        
        # Get font name
        font_name = None
        if paragraph.runs and paragraph.runs[0].font.name:
            font_name = paragraph.runs[0].font.name
        
        # Get text color
        text_color = None
        if paragraph.runs:
            run = paragraph.runs[0]
            try:
                if run.font.color and run.font.color.type == MSO_COLOR_TYPE.RGB:
                    rgb = run.font.color.rgb
                    text_color = (rgb[0], rgb[1], rgb[2])
            except (AttributeError, ValueError):
                # Color not available or not RGB type
                pass
        
        return {
            'text': text,
            'font_size': font_size,
            'font_name': font_name,
            'color': text_color,
            'word_count': len(text.split())
        }
    


    def analyze(self):
        """Run all checks and generate report"""
        self.parse_slides()
        
        report = {
            'title': self.fileName,
            'overall_score': 100,
            'total_slides': len(self.slides_data),
            'estimated_time_minutes': self._estimate_time(),
            'critical_issues': [],
            'warnings': [],
            'strengths': [],
            'slide_details': []
        }
        
        # Global checks
        self._check_font_consistency(report)
        self._check_slide_count(report)
        
        # Per-slide checks
        for slide in self.slides_data:
            slide_report = self._analyze_slide(slide)
            report['slide_details'].append(slide_report)
            
            # Aggregate issues
            for issue in slide_report['critical_issues']:

                val = f"Slide {slide['slide_number']}: {issue}"

                if not (val in report['critical_issues']):

                    report['critical_issues'].append(val)

                
            for issue in slide_report['warnings']:

                val = f"Slide {slide['slide_number']}: {issue}"

                if not (val in report['warnings']):

                    report['warnings'].append(val)
        
        # Calculate final score
        report['overall_score'] -= len(report['critical_issues']) * 5
        report['overall_score'] -= len(report['warnings']) * 2
        report['overall_score'] = max(0, min(100, report['overall_score']))
        
        # Find strengths
        self._identify_strengths(report)
        
        return report
    


    def _analyze_slide(self, slide):
        """Analyze individual slide"""
        result = {
            'slide_number': slide['slide_number'],
            'critical_issues': [],
            'warnings': [],
            'score': 100
        }
        
        # 1. Text size check
        for elem in slide['text_elements']:
            size = elem['font_size']
            if size:
                if size < 18:
                    
                    val = f"Text too small ({size}pt, recommend >18pt)"

                    if not val in result['critical_issues']:

                        result['critical_issues'].append(val)


                elif size < 24 and not slide['has_title']:

                    val = f"Body text could be larger ({size}pt, recommend >24pt for readability)"

                    if not val in result['warnings']:

                        result['warnings'].append(val)
        
        # 2. Word count
        total_words = sum(elem['word_count'] for elem in slide['text_elements'])
        if total_words > 50:
            result['warnings'].append(
                f"{total_words} words (recommend <50 for clarity)"
            )
        
        # 3. Bullet points
        bullet_count = sum(
            1 for elem in slide['text_elements']
            if any(marker in elem['text'] for marker in ['•', '·', '-', '→'])
            or elem['text'].lstrip()[0:2] in ['- ', '• ', '* ']
        )
        if bullet_count > 6:
            result['warnings'].append(
                f"{bullet_count} bullet points (recommend <6)"
            )
        
        # 4. Missing title
        if not slide['has_title'] and slide['text_elements']:
            result['warnings'].append("No title detected")
        
        # 5. Color contrast (if colors available)
        self._check_contrast(slide, result)
        
        # 6. Image-to-text ratio
        if slide['images'] and slide['text_elements']:
            ratio = len(slide['images']) / len(slide['text_elements'])
            if ratio > 2:
                result['warnings'].append(
                    "More images than text (may be cluttered)"
                )
        elif not slide['images'] and not slide['text_elements']:
            result['warnings'].append("Empty slide")


        result['score'] -= len(result['critical_issues']) * 5
        result['score'] -= len(result['warnings']) * 2
        result['score'] = max(0, min(100, result['score']))
        
        return result
    


    def _check_contrast(self, slide, result):
        """Check color contrast using colorsys"""
        # Look for text with colors
        colors_found = [elem['color'] for elem in slide['text_elements'] 
                       if elem['color']]
        
        if colors_found:
            # Assume white background (common default)
            bg_color = (255, 255, 255)
            
            for text_color in colors_found:
                ratio = self._calculate_contrast_ratio(text_color, bg_color)
                if ratio < 4.5:
                    result['critical_issues'].append(
                        f"Poor contrast ratio {ratio:.1f}:1 (needs >4.5:1 for WCAG)"
                    )
                    break  # Only flag once per slide
    


    def _calculate_contrast_ratio(self, rgb1, rgb2):
        """Calculate WCAG contrast ratio using colorsys"""
        def relative_luminance(rgb):
            # Normalize RGB
            r, g, b = [x / 255.0 for x in rgb]
            
            # Convert to linear RGB
            def adjust(c):
                return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
            
            r, g, b = adjust(r), adjust(g), adjust(b)
            
            # Calculate luminance
            return 0.2126 * r + 0.7152 * g + 0.0722 * b
        
        l1 = relative_luminance(rgb1)
        l2 = relative_luminance(rgb2)
        
        lighter = max(l1, l2)
        darker = min(l1, l2)
        
        return (lighter + 0.05) / (darker + 0.05)
    


    def _check_font_consistency(self, report):
        """Check if using too many fonts"""
        fonts_used = set()
        for slide in self.slides_data:
            for elem in slide['text_elements']:
                if elem['font_name']:
                    fonts_used.add(elem['font_name'])
        
        if len(fonts_used) > 2:
            report['warnings'].append(
                f"Using {len(fonts_used)} different fonts (recommend <2 for consistency)"
            )
        elif len(fonts_used) <= 2 and fonts_used:
            report['strengths'].append("Consistent font usage")
    


    def _check_slide_count(self, report):
        """Check if slide count is appropriate"""
        count = len(self.slides_data)
        if count > 20:
            report['warnings'].append(
                f"{count} slides (may be too long for typical presentation)"
            )
        elif count <= 15:
            report['strengths'].append(f"Appropriate slide count ({count})")
    


    def _estimate_time(self):
        """Estimate presentation time in minutes"""
        total_time = 0
        for slide in self.slides_data:
            words = sum(elem['word_count'] for elem in slide['text_elements'])
            # ~150 words per minute speaking rate + time for visuals
            time = (words / 150) + 0.5  # 30 seconds baseline per slide
            total_time += time
        
        return round(total_time, 1)
    


    def _identify_strengths(self, report):
        """Identify positive aspects"""
        # Good image balance
        slides_with_images = sum(1 for s in self.slides_data if s['images'])
        if 0.3 <= slides_with_images / len(self.slides_data) <= 0.7:
            report['strengths'].append("Good image balance")
        
        # Not too wordy
        avg_words = sum(
            sum(e['word_count'] for e in s['text_elements'])
            for s in self.slides_data
        ) / len(self.slides_data)
        
        if avg_words <= 40:
            report['strengths'].append("Concise text (easy to read)")